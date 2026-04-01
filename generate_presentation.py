#!/usr/bin/env python3
"""
Generate thesis_presentation.pptx
Thai Text Segmenter System – thesis committee presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as MSO

# ── Palette ───────────────────────────────────────────────────────────────────
DARK_BLUE   = RGBColor(0x1F, 0x38, 0x64)
MID_BLUE    = RGBColor(0x2E, 0x75, 0xB6)
LIGHT_BLUE  = RGBColor(0xBD, 0xD7, 0xEE)
PALE_BLUE   = RGBColor(0xEB, 0xF3, 0xFB)
GOLD        = RGBColor(0xC5, 0x91, 0x1A)
YELLOW      = RGBColor(0xFF, 0xE6, 0x99)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE   = RGBColor(0xF7, 0xF7, 0xF7)
DARK_GRAY   = RGBColor(0x26, 0x26, 0x26)
MID_GRAY    = RGBColor(0x76, 0x76, 0x76)
GREEN       = RGBColor(0x37, 0x86, 0x2C)
LIGHT_GREEN = RGBColor(0xE2, 0xEF, 0xDA)
ORANGE      = RGBColor(0xED, 0x7D, 0x31)
LIGHT_ORANGE= RGBColor(0xFF, 0xE8, 0xD0)
RED         = RGBColor(0xC0, 0x00, 0x00)
LIGHT_RED   = RGBColor(0xFF, 0xD7, 0xD7)
TEAL        = RGBColor(0x00, 0x70, 0x6A)
LIGHT_TEAL  = RGBColor(0xCC, 0xEE, 0xEC)
PURPLE      = RGBColor(0x7B, 0x2F, 0x8C)
LIGHT_PURPLE= RGBColor(0xF0, 0xDD, 0xF8)

TH = "TH Sarabun New"
EN = "Calibri"

W = 13.333
H = 7.5

# ── Presentation setup ────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(W)
prs.slide_height = Inches(H)
BLANK = prs.slide_layouts[6]

# ── Low-level helpers ─────────────────────────────────────────────────────────

def new_slide():
    return prs.slides.add_slide(BLANK)

def rect(sl, l, t, w, h, fill=None, line=None, lw=1.0, rounded=False):
    st = MSO.ROUNDED_RECTANGLE if rounded else MSO.RECTANGLE
    shp = sl.shapes.add_shape(st, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line:
        shp.line.color.rgb = line
        shp.line.width = Pt(lw)
    else:
        shp.line.fill.background()
    return shp

def label(sl, text, l, t, w, h, size=18, bold=False, italic=False,
          color=DARK_GRAY, align=PP_ALIGN.LEFT, font=TH, wrap=True):
    txb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text          = text
    run.font.size     = Pt(size)
    run.font.bold     = bold
    run.font.italic   = italic
    run.font.color.rgb= color
    run.font.name     = font
    return txb

def mlabel(sl, lines, l, t, w, h, size=16, color=DARK_GRAY,
           align=PP_ALIGN.LEFT, font=TH):
    """Multi-line label: lines = [(text, bold), ...]"""
    txb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, (text, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text           = text
        run.font.size      = Pt(size)
        run.font.bold      = bold
        run.font.color.rgb = color
        run.font.name      = font

def header(sl, title, subtitle=None, chapter=None, bg=DARK_BLUE):
    rect(sl, 0, 0, W, 1.15, fill=bg)
    x = 0.25
    if chapter:
        rect(sl, 0.25, 0.22, 1.5, 0.7, fill=GOLD, rounded=True)
        label(sl, chapter, 0.25, 0.22, 1.5, 0.7, size=15, bold=True,
              color=WHITE, align=PP_ALIGN.CENTER, font=EN)
        x = 1.9
    label(sl, title, x, 0.08, W - x - 0.2, 0.75, size=30, bold=True,
          color=WHITE, font=TH)
    if subtitle:
        label(sl, subtitle, x, 0.78, W - x - 0.2, 0.4, size=17,
              color=LIGHT_BLUE, font=TH)

def tag_box(sl, text, l, t, w, h, bg, fg=WHITE, size=16, bold=True, rounded=True):
    rect(sl, l, t, w, h, fill=bg, rounded=rounded)
    label(sl, text, l, t, w, h, size=size, bold=bold, color=fg,
          align=PP_ALIGN.CENTER, font=TH)

def metric_card(sl, value, label_text, l, t, bg=MID_BLUE, w=2.4, h=1.4):
    rect(sl, l, t, w, h, fill=bg, rounded=True)
    label(sl, value, l, t + 0.1, w, 0.75, size=34, bold=True,
          color=WHITE, align=PP_ALIGN.CENTER, font=EN)
    label(sl, label_text, l, t + 0.82, w, 0.55, size=15, bold=False,
          color=WHITE, align=PP_ALIGN.CENTER, font=TH)

def pipeline_box(sl, stage_no, stage_name, model, example, l, t,
                 w=2.35, h=1.5, bg=MID_BLUE):
    rect(sl, l, t, w, h, fill=bg, rounded=True)
    # stage number circle
    rect(sl, l + 0.1, t + 0.1, 0.45, 0.45,
         fill=WHITE, rounded=True)
    label(sl, stage_no, l + 0.1, t + 0.08, 0.45, 0.45, size=18, bold=True,
          color=bg, align=PP_ALIGN.CENTER, font=EN)
    label(sl, stage_name, l + 0.6, t + 0.1, w - 0.7, 0.45, size=17, bold=True,
          color=WHITE, font=TH)
    label(sl, model, l + 0.1, t + 0.58, w - 0.2, 0.35, size=13, bold=False,
          color=YELLOW, font=EN)
    label(sl, example, l + 0.1, t + 0.93, w - 0.2, 0.5, size=13, bold=False,
          color=WHITE, font=TH, wrap=True)

# ── SLIDE 1  Title ────────────────────────────────────────────────────────────
sl = new_slide()
# Full-slide gradient-feel background
rect(sl, 0, 0, W, H, fill=DARK_BLUE)
rect(sl, 0, 4.8, W, 2.7, fill=RGBColor(0x17, 0x2B, 0x52))
# Gold accent bar
rect(sl, 0, 3.0, W, 0.07, fill=GOLD)

label(sl, "ระบบการแบ่งส่วนข้อความภาษาไทย",
      0.5, 1.0, W - 1, 1.1, size=42, bold=True,
      color=WHITE, align=PP_ALIGN.CENTER, font=TH)
label(sl, "Thai Text Segmenter System",
      0.5, 2.0, W - 1, 0.7, size=28, bold=False,
      color=LIGHT_BLUE, align=PP_ALIGN.CENTER, font=EN)

rect(sl, 3.5, 3.2, 6.3, 0.06, fill=GOLD)

label(sl, "นางสาวยูนิส  เหลียว   รหัส 6511130010",
      0.5, 3.5, W - 1, 0.65, size=22,
      color=WHITE, align=PP_ALIGN.CENTER, font=TH)
label(sl, "ปริญญานิพนธ์ สาขาวิศวกรรมคอมพิวเตอร์และปัญญาประดิษฐ์",
      0.5, 4.1, W - 1, 0.55, size=19,
      color=LIGHT_BLUE, align=PP_ALIGN.CENTER, font=TH)
label(sl, "มหาวิทยาลัยเทคโนโลยีมหานคร  |  ปีการศึกษา 2568",
      0.5, 4.6, W - 1, 0.5, size=18,
      color=LIGHT_BLUE, align=PP_ALIGN.CENTER, font=TH)

# Bottom tag chips
for i, (txt_s, col) in enumerate([
    ("CRF Model", MID_BLUE), ("Viterbi DP", TEAL),
    ("LST20 Corpus", GREEN), ("POS Tagging", PURPLE)
]):
    tag_box(sl, txt_s, 1.8 + i * 2.5, 5.5, 2.1, 0.55,
            bg=col, size=16)

# ── SLIDE 2  Agenda ───────────────────────────────────────────────────────────
sl = new_slide()
rect(sl, 0, 0, W, H, fill=OFF_WHITE)
header(sl, "สารบัญ / Agenda")

chapters = [
    ("บทที่ 1", "บทนำ: ปัญหาและวัตถุประสงค์",          DARK_BLUE),
    ("บทที่ 2", "ทฤษฎี: Pipeline และอัลกอริทึม",       MID_BLUE),
    ("บทที่ 3", "การออกแบบระบบ",                       TEAL),
    ("บทที่ 4", "การทดลอง",                            GREEN),
    ("บทที่ 5", "ผลและสรุป",                           PURPLE),
]
for i, (ch, desc, col) in enumerate(chapters):
    y = 1.4 + i * 1.1
    rect(sl, 0.4, y, 1.5, 0.82, fill=col, rounded=True)
    label(sl, ch, 0.4, y, 1.5, 0.82, size=18, bold=True,
          color=WHITE, align=PP_ALIGN.CENTER, font=EN)
    rect(sl, 2.1, y, 10.8, 0.82, fill=WHITE, line=col, lw=1.5)
    label(sl, desc, 2.3, y, 10.5, 0.82, size=22, bold=False,
          color=DARK_GRAY, font=TH)

# ── SLIDE 3  Why Thai NLP is Hard ────────────────────────────────────────────
sl = new_slide()
rect(sl, 0, 0, W, H, fill=OFF_WHITE)
header(sl, "ทำไมภาษาไทยถึงยาก?", chapter="บทที่ 1")

# English vs Thai comparison
rect(sl, 0.4, 1.3, 5.8, 1.1, fill=LIGHT_GREEN, line=GREEN, lw=1.5)
label(sl, "English  (มีช่องว่าง)", 0.5, 1.3, 5.6, 0.45, size=16,
      bold=True, color=GREEN, font=TH)
label(sl, '"I  love  Thai  food"', 0.5, 1.65, 5.6, 0.55, size=22,
      bold=True, color=DARK_GRAY, align=PP_ALIGN.CENTER, font=EN)

rect(sl, 0.4, 2.6, 5.8, 1.1, fill=LIGHT_RED, line=RED, lw=1.5)
label(sl, "Thai  (ไม่มีช่องว่าง)", 0.5, 2.6, 5.6, 0.45, size=16,
      bold=True, color=RED, font=TH)
label(sl, '"ฉันชอบอาหารไทย"', 0.5, 2.95, 5.6, 0.55, size=26,
      bold=True, color=DARK_GRAY, align=PP_ALIGN.CENTER, font=TH)

# Possible splits
rect(sl, 0.4, 3.9, 5.8, 2.3, fill=PALE_BLUE, line=MID_BLUE, lw=1.5)
label(sl, "ความกำกวม (Ambiguity):", 0.55, 3.9, 5.5, 0.45,
      size=16, bold=True, color=MID_BLUE, font=TH)
mlabel(sl, [
    ("ตา | กลม  = ตา + กลม (eye + round)", False),
    ("ตากลม     = ตากลม (person's name)", False),
    ("ตาก | ลม  = to dry + wind", False),
], 0.55, 4.35, 5.5, 1.7, size=18, color=DARK_GRAY, font=TH)

# 3 challenge boxes
challenges = [
    ("1", "Word Boundary\nAmbiguity",    "ข้อความเดียวอ่านได้หลายแบบ", RED,    LIGHT_RED),
    ("2", "Compound Words",              "คำประสมที่ตีความได้หลายทาง",  ORANGE, LIGHT_ORANGE),
    ("3", "Unknown Words",               "คำใหม่ คำแสลง คำยืมต่างประเทศ",PURPLE, LIGHT_PURPLE),
]
for i, (num, en_t, th_t, col, bg) in enumerate(challenges):
    x = 6.6 + i * 2.22
    rect(sl, x, 1.3, 2.0, 3.5, fill=bg, line=col, lw=2, rounded=True)
    rect(sl, x + 0.7, 1.3, 0.6, 0.6, fill=col, rounded=True)
    label(sl, num, x + 0.7, 1.3, 0.6, 0.6, size=22, bold=True,
          color=WHITE, align=PP_ALIGN.CENTER, font=EN)
    label(sl, en_t, x + 0.1, 2.0, 1.8, 0.85, size=14, bold=True,
          color=col, align=PP_ALIGN.CENTER, font=EN)
    label(sl, th_t, x + 0.1, 2.9, 1.8, 1.7, size=14,
          color=DARK_GRAY, align=PP_ALIGN.CENTER, font=TH)

label(sl, "→ ต้องการระบบอัจฉริยะในการระบุขอบเขตคำ", 0.4, 5.95, 13.0, 0.55,
      size=20, bold=True, color=DARK_BLUE, font=TH)

# ── SLIDE 4  Objectives ───────────────────────────────────────────────────────
sl = new_slide()
rect(sl, 0, 0, W, H, fill=OFF_WHITE)
header(sl, "วัตถุประสงค์และแนวทาง", chapter="บทที่ 1")

label(sl, "พัฒนาระบบ NLP ภาษาไทยแบบ Multi-stage Pipeline",
      0.4, 1.3, 12.8, 0.6, size=24, bold=True, color=DARK_BLUE, font=TH)

objectives = [
    ("🎯", "แก้ปัญหาความกำกวมของขอบเขตคำ (Word Boundary Ambiguity)"),
    ("🎯", "จัดการคำประสม (Compound Words) ได้อย่างถูกต้อง"),
    ("🎯", "รองรับคำนอกพจนานุกรม (Unknown Words)"),
    ("🎯", "ฝึกและประเมินบนคลังข้อมูล LST20 มาตรฐาน"),
]
for i, (icon, obj) in enumerate(objectives):
    y = 2.05 + i * 0.75
    rect(sl, 0.4, y, 12.5, 0.65, fill=PALE_BLUE, line=MID_BLUE, lw=1, rounded=True)
    label(sl, obj, 0.7, y, 12.1, 0.65, size=20, color=DARK_GRAY, font=TH)

# Core model highlight
rect(sl, 0.4, 5.2, 12.5, 1.5, fill=DARK_BLUE, rounded=True)
label(sl, "แบบจำลองหลัก: Conditional Random Field (CRF)",
      0.6, 5.25, 12.1, 0.55, size=22, bold=True, color=GOLD, font=TH)
mlabel(sl, [
    ("ใช้ CRF เดียวกัน 3 ขั้นตอน: MTU Extraction  →  Syllable Identification  →  POS Tagging", False),
    ("ขั้นตอน Word Segmentation ใช้ Viterbi Dynamic Programming + LST20 Dictionary", False),
], 0.6, 5.78, 12.1, 0.85, size=18, color=WHITE, font=TH)

# ── SLIDE 5  System Architecture ─────────────────────────────────────────────
sl = new_slide()
rect(sl, 0, 0, W, H, fill=OFF_WHITE)
header(sl, "สถาปัตยกรรมระบบ (Client–Server)", chapter="บทที่ 2")

# Client box
rect(sl, 0.3, 1.3, 3.8, 5.5, fill=PALE_BLUE, line=MID_BLUE, lw=2, rounded=True)
label(sl, "CLIENT", 0.3, 1.3, 3.8, 0.5, size=18, bold=True,
      color=MID_BLUE, align=PP_ALIGN.CENTER, font=EN)
for i, comp in enumerate(["React\n(UI)", "HTML Mode\n+ Text Mode",
                           "History\n(localStorage)", "Output\nSettings"]):
    x = 0.5 if i % 2 == 0 else 2.2
    y = 2.0 + (i // 2) * 1.65
    rect(sl, x, y, 1.5, 1.3, fill=LIGHT_BLUE, line=MID_BLUE, lw=1, rounded=True)
    label(sl, comp, x, y, 1.5, 1.3, size=15, bold=False,
          color=DARK_BLUE, align=PP_ALIGN.CENTER, font=TH)

# Arrow
label(sl, "HTTP\nPOST / JSON", 4.3, 3.2, 1.5, 1.1, size=14,
      color=MID_GRAY, align=PP_ALIGN.CENTER, font=EN)
label(sl, "⟺", 4.2, 2.9, 1.7, 0.6, size=36, bold=True,
      color=MID_BLUE, align=PP_ALIGN.CENTER, font=EN)

# Server box
rect(sl, 6.0, 1.3, 4.0, 5.5, fill=RGBColor(0xF0, 0xF7, 0xF0),
     line=GREEN, lw=2, rounded=True)
label(sl, "SERVER (FastAPI + Uvicorn)", 6.0, 1.3, 4.0, 0.5,
      size=14, bold=True, color=GREEN, align=PP_ALIGN.CENTER, font=EN)
for i, comp in enumerate(["HTML Parser", "NLP Pipeline\n(CRF + Viterbi)",
                           "Auth API", "SQLite DB"]):
    x = 6.2 if i % 2 == 0 else 7.9
    y = 2.0 + (i // 2) * 1.65
    rect(sl, x, y, 1.6, 1.3, fill=LIGHT_GREEN, line=GREEN, lw=1, rounded=True)
    label(sl, comp, x, y, 1.6, 1.3, size=14, bold=False,
          color=RGBColor(0x20, 0x60, 0x20), align=PP_ALIGN.CENTER, font=EN)

# Models box
rect(sl, 10.2, 1.3, 2.9, 5.5, fill=LIGHT_PURPLE, line=PURPLE, lw=2, rounded=True)
label(sl, "AI Models (.pkl)", 10.2, 1.3, 2.9, 0.5, size=14, bold=True,
      color=PURPLE, align=PP_ALIGN.CENTER, font=EN)
for i, m in enumerate(["MTU CRF", "Syllable CRF", "Viterbi +\nDictionary",
                        "POS CRF"]):
    y = 2.0 + i * 1.2
    rect(sl, 10.35, y, 2.5, 0.95, fill=WHITE, line=PURPLE, lw=1, rounded=True)
    label(sl, m, 10.35, y, 2.5, 0.95, size=14, bold=True,
          color=PURPLE, align=PP_ALIGN.CENTER, font=EN)

label(sl, "React + TypeScript   |   FastAPI + Python 3.13   |   python-crfsuite   |   SQLite3",
      0.3, 6.85, 13.0, 0.45, size=14, color=MID_GRAY,
      align=PP_ALIGN.CENTER, font=EN)

# ── SLIDE 6  4-Stage Pipeline ─────────────────────────────────────────────────
sl = new_slide()
rect(sl, 0, 0, W, H, fill=OFF_WHITE)
header(sl, "4-Stage NLP Pipeline", subtitle="Text → MTU → Syllable → Word → POS",
       chapter="บทที่ 2")

stages = [
    ("1", "MTU\nExtraction",  "CRF\n(char-level)",  "ฉัน|ชอบ|อา|หาร|ไทย",        DARK_BLUE),
    ("2", "Syllable\nID",     "CRF\n(MTU-level)",   "[ฉัน][ชอบ][อา-หาร][ไทย]",   MID_BLUE),
    ("3", "Word\nSegmentation","Viterbi DP\n+ Dict", "ฉัน | ชอบ | อาหาร | ไทย",   TEAL),
    ("4", "POS\nTagging",     "CRF\n(word-level)",  "ฉัน(PA) ชอบ(VV)\nอาหาร(NN) ไทย(AJ)", PURPLE),
]
for i, (num, name, model, ex, col) in enumerate(stages):
    x = 0.3 + i * 3.25
    pipeline_box(sl, num, name, model, ex, x, 1.5, w=2.95, h=4.0, bg=col)
    if i < 3:
        label(sl, "→", x + 3.0, 3.1, 0.4, 0.6, size=30, bold=True,
              color=GOLD, align=PP_ALIGN.CENTER, font=EN)

# Input and output
rect(sl, 0.0, 1.5, 0.32, 4.0, fill=GOLD)
label(sl, "INPUT\nTEXT", 0.0, 2.9, 0.35, 1.2, size=11, bold=True,
      color=WHITE, align=PP_ALIGN.CENTER, font=EN)

# Example text row
rect(sl, 0.3, 5.65, 12.8, 0.65, fill=YELLOW, line=GOLD, lw=1.5, rounded=True)
label(sl, 'ตัวอย่าง: "นักเรียนทุกคนไปโรงพยาบาล"', 0.5, 5.68, 7.0, 0.58,
      size=18, bold=True, color=DARK_BLUE, font=TH)
label(sl, "→  นักเรียน | ทุก | คน | ไป | โรงพยาบาล", 7.5, 5.68, 5.6, 0.58,
      size=18, bold=True, color=GREEN, font=TH)

# ── SLIDE 7  Stage 1: MTU ─────────────────────────────────────────────────────
sl = new_slide()
rect(sl, 0, 0, W, H, fill=OFF_WHITE)
header(sl, "Stage 1: MTU Extraction", subtitle="Minimal Text Unit – smallest building block",
       chapter="บทที่ 2")

# Left: what is MTU
rect(sl, 0.3, 1.3, 5.5, 2.0, fill=PALE_BLUE, line=MID_BLUE, lw=1.5, rounded=True)
label(sl, "MTU คืออะไร?", 0.4, 1.3, 5.2, 0.5, size=18, bold=True,
      color=MID_BLUE, font=TH)
mlabel(sl, [
    ("พยัญชนะ 1 ตัว + สระ + วรรณยุกต์ที่กำกับโดยตรง", False),
    ("ตัวอย่าง: เกา  หลี  ษ  ต  ร", False),
    ("พยัญชนะที่ไม่มีสระ = MTU เดี่ยว: ษ ต ร", False),
], 0.4, 1.82, 5.2, 1.35, size=17, color=DARK_GRAY, font=TH)

# Example word breakdown
rect(sl, 0.3, 3.45, 5.5, 2.5, fill=LIGHT_ORANGE, line=ORANGE, lw=1.5, rounded=True)
label(sl, "ตัวอย่าง: เกษตรกรรม", 0.4, 3.5, 5.2, 0.5, size=18, bold=True,
      color=ORANGE, font=TH)
mtus = ["เก", "ษ", "ต", "ร", "กร", "รม"]
for j, m in enumerate(mtus):
    cx = 0.55 + j * 0.82
    rect(sl, cx, 4.1, 0.7, 0.65, fill=WHITE, line=ORANGE, lw=1.5, rounded=True)
    label(sl, m, cx, 4.1, 0.7, 0.65, size=18, bold=True,
          color=ORANGE, align=PP_ALIGN.CENTER, font=TH)
label(sl, "BMES: B  S  S  S  B  E", 0.4, 4.9, 5.2, 0.45,
      size=16, color=DARK_GRAY, font=EN)
label(sl, "→ Labels สร้างอัตโนมัติจากกฎอักขรวิธี (Silver Labels)", 0.4, 5.35, 5.2, 0.5,
      size=14, italic=True, color=MID_GRAY, font=TH)

# Right: 12 features
rect(sl, 6.1, 1.3, 6.9, 5.4, fill=WHITE, line=MID_BLUE, lw=1.5, rounded=True)
label(sl, "12 Character-Level Features", 6.2, 1.3, 6.6, 0.5, size=18, bold=True,
      color=MID_BLUE, font=EN)
features = [
    ("char",          "อักขระปัจจุบัน"),
    ("char_type",     "ประเภท: C/F/U/L/B/T/K"),
    ("char[±1,±2]",   "Context window ±2 ตัว"),
    ("type[±1,±2]",   "ประเภทใน context window"),
    ("bigram [-1,0]", "คู่อักขระก่อนหน้า+ปัจจุบัน"),
    ("bigram [0,+1]", "คู่อักขระปัจจุบัน+ถัดไป"),
    ("is_consonant",  "เป็นพยัญชนะ?"),
    ("is_vowel",      "เป็นสระ?"),
    ("is_tone",       "เป็นวรรณยุกต์?"),
    ("compound_vowel","เป็นสระประสม?"),
]
for j, (feat, desc) in enumerate(features):
    y = 1.88 + j * 0.46
    rect(sl, 6.2, y, 2.2, 0.38, fill=PALE_BLUE, rounded=True)
    label(sl, feat, 6.25, y, 2.1, 0.38, size=13, bold=True,
          color=DARK_BLUE, font=EN)
    label(sl, desc, 8.5, y, 4.4, 0.38, size=14, color=DARK_GRAY, font=TH)

# Result metric
rect(sl, 0.3, 6.1, 12.8, 0.65, fill=DARK_BLUE, rounded=True)
label(sl, "ผลลัพธ์ MTU CRF:   F1 = 99.80%   |   Precision = 99.82%   |   Recall = 99.78%",
      0.5, 6.13, 12.5, 0.57, size=20, bold=True, color=WHITE,
      align=PP_ALIGN.CENTER, font=EN)

# ── SLIDE 8  Stage 2: Syllable ────────────────────────────────────────────────
sl = new_slide()
rect(sl, 0, 0, W, H, fill=OFF_WHITE)
header(sl, "Stage 2: Syllable Identification",
       subtitle="รวม MTU → พยางค์ด้วย CRF + กฎ Rule-based Signal",
       chapter="บทที่ 2")

# flow diagram
rect(sl, 0.3, 1.3, 5.8, 1.0, fill=LIGHT_BLUE, line=MID_BLUE, lw=1.5, rounded=True)
label(sl, "Input: MTU sequence  [เก][ษ][ต][ร][กร][รม]", 0.5, 1.33, 5.5, 0.9,
      size=17, color=DARK_BLUE, font=TH)
label(sl, "↓  CRF predicts BMES labels", 0.5, 2.4, 5.5, 0.45,
      size=16, bold=True, color=MID_BLUE, font=EN)
rect(sl, 0.3, 2.9, 5.8, 0.8, fill=PALE_BLUE, line=MID_BLUE, lw=1.5, rounded=True)
label(sl, "BMES:  B  S  S  S  B  E  →  [เกษ] [ต] [ร] [กรรม]", 0.5, 2.92, 5.5, 0.72,
      size=17, color=DARK_BLUE, font=TH)
label(sl, "↓  Assemble syllables", 0.5, 3.78, 5.5, 0.4,
      size=16, bold=True, color=MID_BLUE, font=EN)
rect(sl, 0.3, 4.25, 5.8, 0.75, fill=LIGHT_GREEN, line=GREEN, lw=1.5, rounded=True)
label(sl, "Output: [เกษ | ต | ร | กรรม]", 0.5, 4.27, 5.5, 0.67,
      size=18, bold=True, color=GREEN, font=TH)

# Stolen consonant repair
rect(sl, 6.3, 1.3, 6.7, 4.2, fill=LIGHT_RED, line=RED, lw=2, rounded=True)
label(sl, "⚠  Stolen Consonant Repair", 6.4, 1.3, 6.5, 0.55,
      size=17, bold=True, color=RED, font=EN)
mlabel(sl, [
    ("ปัญหา: CRF ดึงพยัญชนะท้ายไปรวมพยางค์ถัดไป", True),
    ("", False),
    ("ตัวอย่าง:  วิตกว่า", False),
    ("  ผิด:  [วิต] + [กว่า]", False),
    ("  ถูก:  [วิตก] + [ว่า]", False),
    ("", False),
    ("กฎ: ถ้า syllable[i] ไม่ใน dict", False),
    ("     แต่ syllable[i]+syllable[i+1][0]", False),
    ("     อยู่ใน dict  AND", False),
    ("     syllable[i+1][1:] อยู่ใน dict", False),
    ("→ ย้ายอักขระแรกของ[i+1] กลับไป[i]", True),
], 6.4, 1.9, 6.5, 3.45, size=15, color=DARK_GRAY, font=TH)

# Weakness note
rect(sl, 0.3, 5.15, 5.8, 1.1, fill=YELLOW, line=GOLD, lw=1.5, rounded=True)
label(sl, "⚠  จุดอ่อน: Over-segmentation", 0.5, 5.18, 5.5, 0.45,
      size=17, bold=True, color=RGBColor(0x80, 0x50, 0x00), font=TH)
mlabel(sl, [
    ("Precision = 61.9%  (ตัดมากเกินไป)", False),
    ("Recall    = 99.9%  (ครอบคลุมดีมาก)", False),
], 0.5, 5.6, 5.5, 0.6, size=16, color=DARK_GRAY, font=EN)

# NFA feature
rect(sl, 6.3, 5.6, 6.7, 0.8, fill=PALE_BLUE, line=MID_BLUE, lw=1.5, rounded=True)
label(sl, "NFA_BOUNDARY feature: เพิ่มสัญญาณจากกฎ orthographic_syllabify เป็น binary feature ให้ CRF",
      6.4, 5.63, 6.5, 0.72, size=15, color=DARK_BLUE, font=TH, wrap=True)

# ── SLIDE 9  Viterbi Word Segmentation ───────────────────────────────────────
sl = new_slide()
rect(sl, 0, 0, W, H, fill=OFF_WHITE)
header(sl, "Stage 3: Word Segmentation (Viterbi DP)",
       subtitle="หาเส้นทางการตัดคำที่ให้คะแนนรวมสูงสุด",
       chapter="บทที่ 2")

# Algorithm description
rect(sl, 0.3, 1.3, 6.2, 2.0, fill=PALE_BLUE, line=MID_BLUE, lw=1.5, rounded=True)
label(sl, "อัลกอริทึม Viterbi DP", 0.5, 1.3, 5.8, 0.5,
      size=18, bold=True, color=MID_BLUE, font=TH)
mlabel(sl, [
    ("1. ไล่จากซ้ายไปขวา ทีละตำแหน่งพยางค์", False),
    ("2. ทดลองรวม 1–8 พยางค์เป็นคำผู้สมัคร", False),
    ("3. เก็บเฉพาะคะแนนสะสมสูงสุดแต่ละตำแหน่ง", False),
    ("4. ย้อนกลับหาเส้นทางที่ดีที่สุด", False),
], 0.5, 1.83, 5.8, 1.38, size=17, color=DARK_GRAY, font=TH)

# DP formula
rect(sl, 0.3, 3.4, 6.2, 1.2, fill=DARK_BLUE, rounded=True)
label(sl, "dp[j] = max  { dp[j-k] + Score(s_{j-k+1}…s_j, k) }",
      0.5, 3.45, 5.8, 0.55, size=18, bold=True,
      color=WHITE, align=PP_ALIGN.CENTER, font=EN)
label(sl, "k = 1…8  (จำกัดที่ 8 พยางค์เพราะคำไทยยาวสุดไม่เกิน 8 พยางค์)",
      0.5, 3.98, 5.9, 0.52, size=14, color=LIGHT_BLUE,
      align=PP_ALIGN.CENTER, font=TH)

# Worked example table
rect(sl, 0.3, 4.7, 6.2, 2.0, fill=WHITE, line=MID_BLUE, lw=1, rounded=True)
label(sl, "ตัวอย่าง: โรงพยาบาล  (3 พยางค์)", 0.5, 4.72, 5.8, 0.45,
      size=16, bold=True, color=MID_BLUE, font=TH)
rows = [
    ("k", "คำผู้สมัคร",      "Score",  "dp[3-k]+Score", "ผล"),
    ("1", "บาล",             "3.60",   "57.39+3.60",    "61.0"),
    ("2", "พยาบาล",          "18.5",   "53.77+18.5",    "72.3"),
    ("3", "โรงพยาบาล ★",    "34.04",  "49.91+34.04",   "83.95 ✓"),
]
for ri, row in enumerate(rows):
    for ci, cell in enumerate(row):
        cx = [0.35, 0.72, 2.35, 3.55, 5.55][ci]
        cw = [0.35, 1.6,  1.15, 1.95, 1.5][ci]
        bg = RGBColor(0xE0, 0xF0, 0xFF) if ri == 0 else (
             LIGHT_GREEN if ri == 4 else WHITE)
        rect(sl, cx, 5.2 + ri * 0.28, cw, 0.27, fill=bg,
             line=RGBColor(0xCC, 0xCC, 0xCC), lw=0.5)
        label(sl, cell, cx, 5.2 + ri * 0.28, cw, 0.27, size=12,
              bold=(ri == 0), color=DARK_GRAY,
              align=PP_ALIGN.CENTER, font=EN if ci == 0 else TH)

# Right: example sentence result
rect(sl, 6.7, 1.3, 6.3, 5.1, fill=WHITE, line=GREEN, lw=2, rounded=True)
label(sl, "ผลลัพธ์ตัวอย่าง:", 6.9, 1.3, 5.9, 0.5,
      size=17, bold=True, color=GREEN, font=TH)
label(sl, '"เมื่อวานนักเรียนทุกคนไปโรงพยาบาล"',
      6.9, 1.85, 5.9, 0.65, size=16, color=DARK_GRAY, font=TH)
label(sl, "↓  Viterbi เลือก:", 6.9, 2.55, 5.9, 0.4, size=15,
      color=MID_BLUE, font=TH)
result_words = [
    ("เมื่อวาน", MID_BLUE), ("นักเรียน", TEAL), ("ทุก", PURPLE),
    ("คน", ORANGE),         ("ไป", GREEN),     ("โรงพยาบาล", RED),
]
for j, (word, col) in enumerate(result_words):
    rx = 6.9 + (j % 3) * 2.0
    ry = 3.05 + (j // 3) * 0.85
    rect(sl, rx, ry, 1.75, 0.7, fill=col, rounded=True)
    label(sl, word, rx, ry, 1.75, 0.7, size=19, bold=True,
          color=WHITE, align=PP_ALIGN.CENTER, font=TH)

rect(sl, 6.7, 4.65, 6.3, 1.75, fill=PALE_BLUE, line=MID_BLUE, lw=1.5, rounded=True)
label(sl, "Compound Bonus CB(k):", 6.9, 4.68, 5.9, 0.45,
      size=16, bold=True, color=MID_BLUE, font=EN)
for j, (k, cb) in enumerate([("k=1", "CB = 0"),
                               ("k=2", "CB = +15"),
                               ("k≥3", "CB = +30")]):
    x = 6.9 + j * 2.0
    label(sl, f"{k}: {cb}", x, 5.15, 1.85, 0.4, size=16,
          bold=(j > 0), color=DARK_BLUE, font=EN)

rect(sl, 0.3, 6.75, 13.0, 0.5, fill=DARK_BLUE, rounded=True)
label(sl, "Word Segmentation F1 = 93.6%",
      0.4, 6.78, 13.0, 0.42, size=20, bold=True,
      color=WHITE, align=PP_ALIGN.CENTER, font=EN)

# ── SLIDE 10  Score Formula ────────────────────────────────────────────────────
sl = new_slide()
rect(sl, 0, 0, W, H, fill=OFF_WHITE)
header(sl, "Score Formula – การให้คะแนนคำผู้สมัคร", chapter="บทที่ 2")

# 4 cases
cases = [
    ("Forced Compound\n(ประเทศไทย, พระนาม...)",
     "Score = 1000", "ชนะเสมอ ไม่มีข้อยกเว้น",
     GOLD, YELLOW),
    ("Dictionary Word\n(w ∈ LST20 dict)",
     "Score = −0.6 + ln(ln(f+1)+1)\n        + 2.5 + CB(k)",
     "−0.6: ค่าปรับต่อขอบเขต\n+2.5: bonus ทุกคำในพจนานุกรม\nCB:   bonus คำประสม",
     MID_BLUE, PALE_BLUE),
    ("Unknown Thai (k=2)\nvalidThai ✓, 2 พยางค์",
     "Score = 7.5",
     "คำประสมใหม่ที่ยังไม่ใน dict",
     TEAL, LIGHT_TEAL),
    ("Other / Invalid\nvalidThai ✗ หรือ k≠2",
     "Score = −1.6  or  −2.6",
     "บังคับหลีกเลี่ยง\nการรวมที่ไม่สมเหตุสมผล",
     RED, LIGHT_RED),
]
for i, (title, formula, desc, col, bg) in enumerate(cases):
    x = 0.3 + (i % 2) * 6.5
    y = 1.3 + (i // 2) * 2.8
    rect(sl, x, y, 6.1, 2.6, fill=bg, line=col, lw=2, rounded=True)
    rect(sl, x, y, 6.1, 0.55, fill=col, rounded=True)
    label(sl, title, x + 0.15, y + 0.02, 5.8, 0.5, size=16, bold=True,
          color=WHITE, font=TH)
    label(sl, formula, x + 0.15, y + 0.62, 5.8, 0.9, size=17, bold=True,
          color=col, font=EN)
    label(sl, desc, x + 0.15, y + 1.55, 5.8, 0.95, size=14,
          color=DARK_GRAY, font=TH)

rect(sl, 0.3, 6.75, 13.0, 0.5, fill=DARK_BLUE, rounded=True)
label(sl, "CB(k) = min(k−1, 2) × 15    →    1-syl: +0    2-syl: +15    3+-syl: +30",
      0.3, 6.78, 13.0, 0.44, size=18, bold=True,
      color=WHITE, align=PP_ALIGN.CENTER, font=EN)

# ── SLIDE 11  k-best Reranking ────────────────────────────────────────────────
sl = new_slide()
rect(sl, 0, 0, W, H, fill=OFF_WHITE)
header(sl, "k-best Reranking ด้วย POS CRF",
       subtitle="ทำไม Viterbi อย่างเดียวยังไม่พอ?",
       chapter="บทที่ 2")

# Problem
rect(sl, 0.3, 1.3, 5.8, 1.6, fill=LIGHT_RED, line=RED, lw=1.5, rounded=True)
label(sl, "ปัญหาของ Viterbi 1-best:", 0.5, 1.3, 5.5, 0.45,
      size=18, bold=True, color=RED, font=TH)
mlabel(sl, [
    ("ประเมินคำทีละคำ ไม่รู้บริบทรอบข้าง", False),
    ("เส้นทางที่คะแนนสูงสุดอาจไม่ถูกต้องทางไวยากรณ์", False),
], 0.5, 1.78, 5.5, 1.0, size=17, color=DARK_GRAY, font=TH)

# Solution
rect(sl, 0.3, 3.0, 5.8, 2.0, fill=LIGHT_GREEN, line=GREEN, lw=1.5, rounded=True)
label(sl, "วิธีแก้: k-best + POS Reranking", 0.5, 3.02, 5.5, 0.48,
      size=18, bold=True, color=GREEN, font=TH)
mlabel(sl, [
    ("1. Viterbi เก็บ 5 เส้นทางที่ดีที่สุด", False),
    ("2. ส่งลำดับคำแต่ละเส้นทางเข้า POS CRF", False),
    ("3. CRF ให้ P(tag | word) ต่อคำ", False),
    ("4. คะแนนรวม = Viterbi + λ × POS", False),
    ("5. เลือกเส้นทางที่คะแนนรวมสูงสุด", False),
], 0.5, 3.55, 5.5, 1.35, size=16, color=DARK_GRAY, font=TH)

# Formula
rect(sl, 0.3, 5.1, 5.8, 0.85, fill=DARK_BLUE, rounded=True)
label(sl, "S* = argmax_i  ( V_i  +  λ · P_i )      λ = 0.3",
      0.5, 5.15, 5.5, 0.72, size=18, bold=True,
      color=WHITE, align=PP_ALIGN.CENTER, font=EN)

# POS confidence explanation
rect(sl, 6.3, 1.3, 6.8, 4.7, fill=WHITE, line=PURPLE, lw=1.5, rounded=True)
label(sl, "POS Confidence Score (Pᵢ)", 6.4, 1.3, 6.5, 0.5,
      size=18, bold=True, color=PURPLE, font=EN)
label(sl, "Pᵢ = Σ ln( max_tag P(tag_t | w_t) + ε )",
      6.4, 1.88, 6.5, 0.55, size=16, bold=True,
      color=PURPLE, font=EN)
label(sl, "• ln ผลลัพธ์เป็นลบเสมอ (0 < P ≤ 1)\n• ค่าใกล้ 0 = CRF มั่นใจสูง\n• ค่าห่างจาก 0 = CRF สับสน",
      6.4, 2.5, 6.5, 0.85, size=15, color=DARK_GRAY, font=TH, wrap=True)

# Comparison table
rows2 = [
    ("Path", "ลำดับคำ",               "V_i",  "P_i",  "รวม"),
    ("1 ✓",  "เมื่อวาน|นักเรียน|...", "83.95","-0.40","83.55"),
    ("2",    "เมื่อ|วาน|นักเรียน|...", "72.10","-2.80","71.26"),
    ("3",    "เมื่อวาน|นัก|เรียน|...", "74.50","-2.10","73.87"),
]
for ri, row in enumerate(rows2):
    for ci, cell in enumerate(row):
        cx = [6.35, 6.75, 9.6, 10.7, 11.8][ci]
        cw = [0.38, 2.82, 1.08, 1.08, 1.3][ci]
        bg = RGBColor(0xCC, 0xE5, 0xFF) if ri == 0 else (
             LIGHT_GREEN if ri == 1 else WHITE)
        rect(sl, cx, 3.45 + ri * 0.55, cw, 0.5, fill=bg,
             line=RGBColor(0xCC, 0xCC, 0xCC), lw=0.5)
        label(sl, cell, cx, 3.45 + ri * 0.55, cw, 0.5, size=12,
              bold=(ri == 0 or (ri == 1 and ci == 0)),
              color=DARK_GRAY, align=PP_ALIGN.CENTER, font=TH)

# Shortcut note
rect(sl, 6.3, 5.75, 6.8, 0.55, fill=YELLOW, line=GOLD, lw=1, rounded=True)
label(sl, "⚡  ถ้า gap > 5.0 → ใช้ Viterbi โดยตรง (ไม่รัน POS เพิ่ม)",
      6.4, 5.77, 6.6, 0.48, size=14, color=DARK_GRAY, font=TH)

# ── SLIDE 12  POS Tagging ─────────────────────────────────────────────────────
sl = new_slide()
rect(sl, 0, 0, W, H, fill=OFF_WHITE)
header(sl, "Stage 4: POS Tagging – 16 Tags (LST20 Standard)", chapter="บทที่ 2")

tags = [
    ("NN","คำนาม","บ้าน รถ ครู",DARK_BLUE),
    ("VV","คำกริยา","กิน วิ่ง เดิน",MID_BLUE),
    ("AJ","คำคุณศัพท์","สวย ใหญ่ ดี",TEAL),
    ("AV","คำวิเศษณ์","เร็ว ช้า มาก",GREEN),
    ("AX","คำช่วย","จะ ได้ ให้",PURPLE),
    ("CC","คำเชื่อม","และ หรือ แต่",ORANGE),
    ("CL","ลักษณนาม","ตัว คน อัน",RED),
    ("FX","คำเติม","การ ความ นัก",RGBColor(0x80,0x00,0x80)),
    ("IJ","คำอุทาน","โอ้ อา โธ่",RGBColor(0x00,0x80,0x80)),
    ("NG","คำปฏิเสธ","ไม่ ไม่ได้",RED),
    ("NU","คำบอกจำนวน","หนึ่ง สอง สาม",MID_BLUE),
    ("PA","คำบอกอัตภาพ","ผม ฉัน เขา",TEAL),
    ("PR","คำสรรพนาม","นี้ นั้น ที่",GREEN),
    ("PS","คำบุพบท","ใน บน ของ",PURPLE),
    ("PU","วรรคตอน",", . ? !",DARK_GRAY),
    ("XX","จำแนกไม่ได้","—",MID_GRAY),
]
for i, (tag, name, ex, col) in enumerate(tags):
    r, c = divmod(i, 4)
    x = 0.3 + c * 3.26
    y = 1.3 + r * 1.38
    rect(sl, x, y, 3.0, 1.2, fill=WHITE, line=col, lw=2, rounded=True)
    rect(sl, x, y, 0.72, 1.2, fill=col, rounded=True)
    label(sl, tag, x, y, 0.72, 1.2, size=20, bold=True,
          color=WHITE, align=PP_ALIGN.CENTER, font=EN)
    label(sl, name, x + 0.77, y + 0.05, 2.1, 0.55, size=16, bold=True,
          color=col, font=TH)
    label(sl, ex, x + 0.77, y + 0.58, 2.1, 0.5, size=13,
          color=MID_GRAY, font=TH)

rect(sl, 0.3, 6.83, 13.0, 0.45, fill=DARK_BLUE, rounded=True)
label(sl, "POS Accuracy = 94.27%   |   NG F1 = 99.85%   |   FX F1 = 99.36%   |   CL F1 = 57.39% (ต่ำสุด)",
      0.3, 6.86, 13.0, 0.38, size=16, bold=True,
      color=WHITE, align=PP_ALIGN.CENTER, font=EN)

# ── SLIDE 13  UI Design ────────────────────────────────────────────────────────
sl = new_slide()
rect(sl, 0, 0, W, H, fill=OFF_WHITE)
header(sl, "การออกแบบระบบ – User Interface", chapter="บทที่ 3")

screens = [
    ("🔐 Login", "ตรวจสอบ username/password\nล็อกอัตโนมัติถ้าผิด 3 ครั้ง",DARK_BLUE),
    ("📄 HTML Mode","ลาก-วางไฟล์ .html\nดึงเฉพาะ text nodes",MID_BLUE),
    ("✏ Text Mode", "พิมพ์/วางข้อความ\nปุ่ม Copy ผลลัพธ์",TEAL),
    ("⚙ Settings","เลือก HTML tag (<span>)\nกำหนด CSS class",GREEN),
    ("📜 History","History Card: ชื่อไฟล์ + เวลา + จำนวน Segments\nSearch / Download / Delete",PURPLE),
    ("📊 Statistics","ภาพรวมผล MTU / Syllable / Word / POS\nบน LST20 test set",ORANGE),
]
for i, (name, desc, col) in enumerate(screens):
    r, c = divmod(i, 3)
    x = 0.3 + c * 4.35
    y = 1.3 + r * 2.6
    rect(sl, x, y, 4.0, 2.35, fill=WHITE, line=col, lw=2, rounded=True)
    rect(sl, x, y, 4.0, 0.52, fill=col, rounded=True)
    label(sl, name, x + 0.1, y + 0.04, 3.8, 0.44, size=17, bold=True,
          color=WHITE, font=TH)
    label(sl, desc, x + 0.15, y + 0.62, 3.7, 1.6, size=15,
          color=DARK_GRAY, font=TH)

# ── SLIDE 14  Database Design ─────────────────────────────────────────────────
sl = new_slide()
rect(sl, 0, 0, W, H, fill=OFF_WHITE)
header(sl, "การออกแบบฐานข้อมูล (SQLite3)", chapter="บทที่ 3")

tables = [
    ("roles", DARK_BLUE, [
        ("id",         "TEXT",  "PK"),
        ("name",       "TEXT",  "UNIQUE NOT NULL"),
        ("description","TEXT",  ""),
        ("is_system",  "INT",   "DEFAULT 0"),
    ]),
    ("users", MID_BLUE, [
        ("id",             "TEXT",      "PK"),
        ("username",       "TEXT",      "UNIQUE NOT NULL"),
        ("password_hash",  "TEXT",      "NOT NULL (bcrypt)"),
        ("role_id",        "TEXT",      "FK → roles.id"),
        ("name",           "TEXT",      "NOT NULL"),
        ("login_attempts", "INT",       "DEFAULT 0"),
        ("is_locked",      "INT",       "DEFAULT 0"),
    ]),
    ("sessions", TEAL, [
        ("token",      "TEXT",      "PK"),
        ("user_id",    "TEXT",      "FK → users.id"),
        ("created_at", "TIMESTAMP", "CURRENT_TIMESTAMP"),
        ("expires_at", "TIMESTAMP", "NOT NULL (+30 min)"),
    ]),
]
for i, (tname, col, cols) in enumerate(tables):
    x = 0.3 + i * 4.3
    rect(sl, x, 1.3, 4.0, 0.55, fill=col, rounded=True)
    label(sl, tname, x, 1.3, 4.0, 0.55, size=20, bold=True,
          color=WHITE, align=PP_ALIGN.CENTER, font=EN)
    for j, (cname, ctype, constraint) in enumerate(cols):
        y = 1.9 + j * 0.52
        bg = PALE_BLUE if col == DARK_BLUE else (
             RGBColor(0xE8, 0xF4, 0xFD) if col == MID_BLUE else LIGHT_TEAL)
        rect(sl, x, y, 4.0, 0.47, fill=bg,
             line=RGBColor(0xCC, 0xCC, 0xCC), lw=0.5)
        label(sl, cname, x + 0.1, y, 1.4, 0.47, size=13, bold=True,
              color=DARK_GRAY, font=EN)
        label(sl, ctype, x + 1.5, y, 0.85, 0.47, size=12,
              color=MID_GRAY, font=EN)
        label(sl, constraint, x + 2.35, y, 1.6, 0.47, size=11,
              color=col, font=EN)

# Security notes
rect(sl, 0.3, 6.3, 12.8, 0.9, fill=YELLOW, line=GOLD, lw=1.5, rounded=True)
label(sl, "🔒  Security:  รหัสผ่าน bcrypt  |  Lock บัญชีถ้า login ผิด 3 ครั้ง  |  Session หมดอายุ 30 นาที  |  is_system=1 ป้องกันลบ admin",
      0.5, 6.35, 12.5, 0.78, size=16, color=DARK_GRAY, font=TH)

# ── SLIDE 15  Experiment Setup ────────────────────────────────────────────────
sl = new_slide()
rect(sl, 0, 0, W, H, fill=OFF_WHITE)
header(sl, "การทดลอง – Dataset & Setup", chapter="บทที่ 4")

# Dataset
rect(sl, 0.3, 1.3, 6.0, 4.8, fill=PALE_BLUE, line=MID_BLUE, lw=2, rounded=True)
label(sl, "คลังข้อมูล: LST20", 0.5, 1.3, 5.6, 0.55,
      size=22, bold=True, color=MID_BLUE, font=TH)
stats = [
    ("ประโยคทั้งหมด", "33,687"),
    ("ชุดฝึก (train)", "~26,000"),
    ("ชุดพัฒนา (dev)", "~4,000"),
    ("ชุดทดสอบ (test)", "~3,687"),
    ("โดเมน", "ข่าว / นิตยสาร / วิชาการ / ทั่วไป"),
    ("ป้ายกำกับ", "คำ + POS tag (Gold Labels)"),
    ("MTU/Syllable", "Silver Labels (กฎอักขรวิธี)"),
]
for j, (k, v) in enumerate(stats):
    y = 1.95 + j * 0.59
    rect(sl, 0.45, y, 5.7, 0.52, fill=WHITE if j % 2 == 0 else PALE_BLUE,
         line=LIGHT_BLUE, lw=0.5)
    label(sl, k, 0.55, y, 2.8, 0.52, size=16, bold=True, color=DARK_BLUE, font=TH)
    label(sl, v, 3.4, y, 2.8, 0.52, size=16, color=DARK_GRAY, font=TH)

# Tech stack
rect(sl, 6.6, 1.3, 6.4, 4.8, fill=WHITE, line=GREEN, lw=2, rounded=True)
label(sl, "สภาพแวดล้อม", 6.8, 1.3, 6.0, 0.55,
      size=22, bold=True, color=GREEN, font=TH)
tech = [
    ("OS",        "Windows"),
    ("Python",    "3.13"),
    ("CRF lib",   "python-crfsuite"),
    ("Backend",   "FastAPI + Uvicorn"),
    ("Frontend",  "React + TypeScript"),
    ("Database",  "SQLite3"),
    ("Algorithm", "sklearn-crfsuite (L-BFGS)"),
]
for j, (k, v) in enumerate(tech):
    y = 1.95 + j * 0.59
    rect(sl, 6.7, y, 6.1, 0.52, fill=LIGHT_GREEN if j % 2 == 0 else WHITE,
         line=RGBColor(0xCC, 0xEE, 0xCC), lw=0.5)
    label(sl, k, 6.8, y, 1.8, 0.52, size=16, bold=True, color=GREEN, font=EN)
    label(sl, v, 8.6, y, 4.1, 0.52, size=16, color=DARK_GRAY, font=EN)

# Metrics definition
rect(sl, 0.3, 6.2, 12.8, 1.05, fill=DARK_BLUE, rounded=True)
mlabel(sl, [
    ("ตัวชี้วัด:   Precision = TP/(TP+FP)     Recall = TP/(TP+FN)     F1 = 2·P·R/(P+R)     Accuracy = correct/total", False),
], 0.5, 6.23, 12.5, 0.98, size=17, color=WHITE,
   align=PP_ALIGN.CENTER, font=EN)

# ── SLIDE 16  Web Wrapping Evaluation ─────────────────────────────────────────
sl = new_slide()
rect(sl, 0, 0, W, H, fill=OFF_WHITE)
header(sl, "Web Text Wrapping Evaluation", chapter="บทที่ 4")

label(sl, "ทดสอบกับหน้าเว็บภาษาไทย: นับจุดที่ข้อความตัดผิด ก่อน/หลัง ใช้ระบบ",
      0.3, 1.3, 12.8, 0.55, size=20, color=DARK_GRAY, font=TH)

# 3 viewport boxes
viewports = [
    ("S – Smartphone", "375 px", "Mid-word breaks\nText overflow"),
    ("M – Tablet",     "768 px", "Awkward breaks\nbetween Thai words"),
    ("L – Desktop",    "1280 px","Occasional overflow\nin narrow columns"),
]
for i, (name, px, issues) in enumerate(viewports):
    x = 0.3 + i * 4.35
    rect(sl, x, 1.95, 4.0, 2.8, fill=WHITE, line=MID_BLUE, lw=1.5, rounded=True)
    rect(sl, x, 1.95, 4.0, 0.52, fill=MID_BLUE, rounded=True)
    label(sl, name, x + 0.1, 1.95, 3.8, 0.52, size=17, bold=True,
          color=WHITE, font=EN)
    label(sl, px, x + 0.15, 2.53, 3.7, 0.42, size=16, bold=True,
          color=MID_BLUE, font=EN)
    label(sl, issues, x + 0.15, 3.0, 3.7, 1.6, size=15,
          color=DARK_GRAY, font=EN)

# Formula
rect(sl, 0.3, 4.85, 12.8, 1.05, fill=PALE_BLUE, line=MID_BLUE, lw=1.5, rounded=True)
label(sl, "Improvement (%) = (Before − After) / Before × 100",
      0.5, 4.88, 12.4, 0.48, size=20, bold=True,
      color=DARK_BLUE, align=PP_ALIGN.CENTER, font=EN)
label(sl, "Overall Improvement = (B_S+B_M+B_L − A_S−A_M−A_L) / (B_S+B_M+B_L) × 100",
      0.5, 5.32, 12.4, 0.48, size=17, color=MID_BLUE,
      align=PP_ALIGN.CENTER, font=EN)

# Target
rect(sl, 0.3, 6.0, 12.8, 0.75, fill=GOLD, rounded=True)
label(sl, "🎯  เป้าหมาย: Overall Improvement ≥ 70%",
      0.5, 6.03, 12.4, 0.68, size=22, bold=True,
      color=WHITE, align=PP_ALIGN.CENTER, font=TH)

# ── SLIDE 17  Results Overview ────────────────────────────────────────────────
sl = new_slide()
rect(sl, 0, 0, W, H, fill=OFF_WHITE)
header(sl, "ผลการทดลอง – ประสิทธิภาพทุกขั้นตอน", chapter="บทที่ 5")

# Big metric cards
metrics = [
    ("99.8%", "MTU CRF\nF1 Score",          GREEN,    0.4),
    ("61.9%", "Syllable CRF\nPrecision",     ORANGE,   3.0),
    ("93.6%", "Word Seg\nF1 Score",          MID_BLUE, 5.6),
    ("94.3%", "POS Tagging\nAccuracy",       PURPLE,   8.2),
]
for val, lbl, col, x in metrics:
    metric_card(sl, val, lbl, x, 1.3, bg=col, w=2.4, h=1.8)

# Detailed table
headers = ["ขั้นตอน", "ประเภทป้ายกำกับ", "Precision", "Recall", "F1 / Acc"]
col_w   = [3.0, 3.0, 1.7, 1.7, 2.0]
col_x   = [0.4, 3.4, 6.4, 8.1, 9.8]

for ci, (h, x, w) in enumerate(zip(headers, col_x, col_w)):
    rect(sl, x, 3.25, w, 0.48, fill=DARK_BLUE)
    label(sl, h, x, 3.25, w, 0.48, size=15, bold=True,
          color=WHITE, align=PP_ALIGN.CENTER, font=TH)

rows3 = [
    ("MTU CRF",         "Silver (กฎอักขรวิธี)", "99.82%", "99.78%", "99.80%", GREEN),
    ("Syllable CRF",    "Silver (กฎอักขรวิธี)", "61.92%", "99.90%", "76.44%", ORANGE),
    ("Word Segmentation","Gold (LST20)",          "93.5%",  "93.7%",  "93.6%",  MID_BLUE),
    ("POS Tagging",     "Gold (LST20)",          "—",      "—",      "94.27%", PURPLE),
]
for ri, (stage, label_type, p, r, f, col) in enumerate(rows3):
    y = 3.78 + ri * 0.68
    values = [stage, label_type, p, r, f]
    for ci, (val, x, w) in enumerate(zip(values, col_x, col_w)):
        bg = LIGHT_GREEN if col == GREEN else (
             LIGHT_ORANGE if col == ORANGE else (
             PALE_BLUE if col == MID_BLUE else LIGHT_PURPLE))
        rect(sl, x, y, w, 0.62,
             fill=bg if ci == 0 else (WHITE if ri % 2 == 0 else OFF_WHITE),
             line=RGBColor(0xCC, 0xCC, 0xCC), lw=0.5)
        b = ci == 0
        label(sl, val, x + 0.05, y, w - 0.05, 0.62, size=15,
              bold=b, color=col if ci == 0 else DARK_GRAY,
              align=PP_ALIGN.CENTER if ci >= 2 else PP_ALIGN.LEFT,
              font=TH)

label(sl, "ทดสอบบน LST20 test set  (33,687 ประโยค)",
      0.4, 6.65, 12.8, 0.45, size=16, color=MID_GRAY,
      align=PP_ALIGN.CENTER, font=TH)

# ── SLIDE 18  Strengths ────────────────────────────────────────────────────────
sl = new_slide()
rect(sl, 0, 0, W, H, fill=OFF_WHITE)
header(sl, "จุดแข็ง (Strengths)", chapter="บทที่ 5")

strengths = [
    ("MTU CRF F1 = 99.80%",
     "กฎอักขรวิธีที่ออกแบบให้เหมาะกับโครงสร้างไทย + Feature 12 ประเภท\nทำให้ระบบระบุขอบเขต MTU ได้แม่นยำสูงมาก"),
    ("Compound Word Detection",
     "Compound Bonus CB(k) ใน Score function บังคับให้ Viterbi\nชอบรวมพยางค์เป็นคำยาว เช่น โรงพยาบาล (+30 bonus)"),
    ("POS Accuracy = 94.27%",
     "NG (ปฏิเสธ) F1 = 99.85%   FX (คำเติม) F1 = 99.36%\nคำที่มีรูปแบบตายตัวจำแนกได้แม่นยำมาก"),
    ("k-best + POS Reranking",
     "ใช้บริบทระดับประโยค (POS) ช่วยแก้กรณีที่ Viterbi ตัดคำผิด\nเช่น นักเรียน vs. นัก+เรียน"),
    ("Stolen Consonant Repair",
     "กลไกซ่อมขอบเขตพยางค์อัตโนมัติก่อน Viterbi\nลดข้อผิดพลาดที่ส่งต่อจาก Syllable CRF"),
]
for i, (title, desc) in enumerate(strengths):
    y = 1.3 + i * 1.07
    rect(sl, 0.3, y, 0.55, 0.88, fill=GREEN, rounded=True)
    label(sl, "✓", 0.3, y, 0.55, 0.88, size=22, bold=True,
          color=WHITE, align=PP_ALIGN.CENTER, font=EN)
    rect(sl, 0.95, y, 12.1, 0.88, fill=LIGHT_GREEN, line=GREEN, lw=1, rounded=True)
    label(sl, title, 1.1, y + 0.02, 12.0, 0.38, size=18, bold=True,
          color=GREEN, font=TH)
    label(sl, desc, 1.1, y + 0.4, 11.8, 0.45, size=14,
          color=DARK_GRAY, font=TH)

# ── SLIDE 19  Weaknesses ───────────────────────────────────────────────────────
sl = new_slide()
rect(sl, 0, 0, W, H, fill=OFF_WHITE)
header(sl, "จุดอ่อนและข้อจำกัด (Weaknesses & Limitations)", chapter="บทที่ 5")

weaknesses = [
    ("Syllable CRF Over-segmentation",
     "Precision = 61.9%  →  ตัดพยางค์มากเกินไป\nส่งข้อผิดพลาดไปยังขั้นตอน Word Segmentation", RED),
    ("CL (ลักษณนาม) F1 = 57.39%",
     "ลักษณนามทำหน้าที่ได้หลายอย่าง (นาม/กริยา) ทำให้ CRF สับสน\nตัวอย่างน้อยในบางบริบท", ORANGE),
    ("Silver Labels (MTU & Syllable)",
     "ป้ายกำกับ MTU และ Syllable สร้างจากกฎอัตโนมัติ ไม่ใช่ผู้เชี่ยวชาญ\nตัวเลข F1 สะท้อนความสอดคล้องกับกฎ ≠ ความถูกต้องภาษาศาสตร์", PURPLE),
    ("In-domain Bias",
     "พจนานุกรม Viterbi สร้างจาก LST20 training set เดียวกัน\nคะแนน Word Seg อาจสูงกว่าความเป็นจริงบน corpus นอก LST20", MID_BLUE),
    ("Unknown Words / คำใหม่",
     "คำแสลง คำยืม คำเฉพาะทางที่ไม่ใน LST20\nViterbi อาจตัดผิดเพราะไม่มีในพจนานุกรม", DARK_GRAY),
]
for i, (title, desc, col) in enumerate(weaknesses):
    y = 1.3 + i * 1.06
    rect(sl, 0.3, y, 0.55, 0.88, fill=col, rounded=True)
    label(sl, "⚠", 0.3, y, 0.55, 0.88, size=20, bold=True,
          color=WHITE, align=PP_ALIGN.CENTER, font=EN)
    rect(sl, 0.95, y, 12.1, 0.88, fill=OFF_WHITE, line=col, lw=1.5, rounded=True)
    label(sl, title, 1.1, y + 0.02, 12.0, 0.38, size=18, bold=True,
          color=col, font=TH)
    label(sl, desc, 1.1, y + 0.4, 11.8, 0.45, size=14,
          color=DARK_GRAY, font=TH)

# ── SLIDE 20  Conclusion ──────────────────────────────────────────────────────
sl = new_slide()
rect(sl, 0, 0, W, H, fill=DARK_BLUE)
rect(sl, 0, 5.5, W, 2.0, fill=RGBColor(0x17, 0x2B, 0x52))
rect(sl, 0, 3.1, W, 0.06, fill=GOLD)

label(sl, "สรุป (Conclusion)", 0.5, 0.4, W - 1, 0.75, size=36, bold=True,
      color=WHITE, align=PP_ALIGN.CENTER, font=TH)

conclusions = [
    ("🔬", "พัฒนาระบบ 4-stage NLP Pipeline สำหรับตัดคำและ POS tagging ภาษาไทย"),
    ("🧠", "ใช้ CRF เป็นแกนหลัก 3 ขั้นตอน + Viterbi DP สำหรับ Word Segmentation"),
    ("📊", "Word F1 = 93.6%  |  POS Accuracy = 94.3%  บนคลัง LST20"),
    ("🌐", "นำไปใช้งานจริงในระบบ Web Text Wrapping ผ่าน React + FastAPI"),
]
for i, (icon, text) in enumerate(conclusions):
    y = 1.3 + i * 0.95
    rect(sl, 0.5, y, 12.3, 0.82, fill=RGBColor(0x2A, 0x4A, 0x7F), rounded=True)
    label(sl, icon + "  " + text, 0.7, y, 12.0, 0.82, size=20,
          color=WHITE, font=TH)

label(sl, "Future Work: เพิ่ม Gold Syllable Labels · ปรับปรุง CL tagging · ทดสอบ Out-of-domain",
      0.5, 5.6, W - 1, 0.55, size=17, color=LIGHT_BLUE,
      align=PP_ALIGN.CENTER, font=TH)
label(sl, "ขอขอบพระคุณ ผศ.ดร.ณัฐพงษ์ จันทร์แดง อาจารย์ที่ปรึกษา",
      0.5, 6.15, W - 1, 0.5, size=18, bold=True,
      color=GOLD, align=PP_ALIGN.CENTER, font=TH)

# ── SLIDE 21  Q&A ─────────────────────────────────────────────────────────────
sl = new_slide()
rect(sl, 0, 0, W, H, fill=DARK_BLUE)
rect(sl, 0, 0, W, 0.12, fill=GOLD)
rect(sl, 0, H - 0.12, W, 0.12, fill=GOLD)

label(sl, "ขอบคุณครับ / ค่ะ", 0.5, 1.5, W - 1, 1.0, size=52, bold=True,
      color=WHITE, align=PP_ALIGN.CENTER, font=TH)
label(sl, "Thank You", 0.5, 2.5, W - 1, 0.8, size=36,
      color=LIGHT_BLUE, align=PP_ALIGN.CENTER, font=EN)

rect(sl, 3.5, 3.5, 6.3, 0.06, fill=GOLD)

label(sl, "Q & A", 0.5, 3.7, W - 1, 0.9, size=48, bold=True,
      color=GOLD, align=PP_ALIGN.CENTER, font=EN)

# Summary metrics row
for i, (val, lbl, col) in enumerate([
    ("99.8%", "MTU F1",       GREEN),
    ("93.6%", "Word F1",      MID_BLUE),
    ("94.3%", "POS Accuracy", PURPLE),
    ("4-stage", "Pipeline",   TEAL),
]):
    metric_card(sl, val, lbl, 1.3 + i * 2.7, 4.8, bg=col, w=2.4, h=1.5)

label(sl, "นางสาวยูนิส เหลียว  |  6511130010  |  วิศวกรรมคอมพิวเตอร์และปัญญาประดิษฐ์  |  MUT 2568",
      0.5, 6.7, W - 1, 0.5, size=15, color=LIGHT_BLUE,
      align=PP_ALIGN.CENTER, font=TH)

# ── Save ──────────────────────────────────────────────────────────────────────
out = "thesis_presentation.pptx"
prs.save(out)
print(f"Saved: {out}  ({len(prs.slides)} slides)")
